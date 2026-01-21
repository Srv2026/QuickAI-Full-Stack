"""
Python script to convert QuickAI_Presentation.md to PowerPoint (.pptx)
Requires: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import re

def parse_markdown_slides(md_file):
    """Parse markdown file and extract slides"""
    with open(md_file, 'r', encoding='utf-8') as f:
        content = f.read()
    
    # Split by slide separators (---)
    slides = []
    current_slide = []
    
    for line in content.split('\n'):
        if line.strip() == '---':
            if current_slide:
                slides.append('\n'.join(current_slide))
                current_slide = []
        else:
            current_slide.append(line)
    
    if current_slide:
        slides.append('\n'.join(current_slide))
    
    return slides

def parse_slide_content(slide_text):
    """Parse slide content to extract title and body"""
    lines = slide_text.strip().split('\n')
    title = ""
    body = []
    in_code_block = False
    
    for line in lines:
        # Skip empty lines at start
        if not line.strip() and not title:
            continue
        
        # Check for code blocks
        if line.strip().startswith('```'):
            in_code_block = not in_code_block
            continue
        
        if in_code_block:
            body.append(line)
            continue
        
        # Extract title (first # heading or Slide X: Title)
        if line.startswith('# ') and not title:
            title = line.replace('# ', '').strip()
            # Remove "Slide X:" prefix if present
            title = re.sub(r'^Slide \d+:\s*', '', title)
        elif line.startswith('## ') and not title:
            title = line.replace('## ', '').strip()
            title = re.sub(r'^Slide \d+:\s*', '', title)
        elif 'Slide' in line and ':' in line and not title:
            # Extract title from "Slide X: Title" format
            title = line.split(':', 1)[1].strip() if ':' in line else line.strip()
        elif line.strip() and not title:
            # First non-empty line becomes title
            title = line.strip()
        else:
            # Body content
            if line.strip():
                body.append(line)
    
    return title, '\n'.join(body)

def add_slide(prs, title, content):
    """Add a slide to presentation"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout
    
    # Set title
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Format title
    title_frame = title_shape.text_frame
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Set content
    if content.strip():
        content_shape = slide.placeholders[1]
        tf = content_shape.text_frame
        tf.text = content
        
        # Format content
        for paragraph in tf.paragraphs:
            paragraph.font.size = Pt(18)
            paragraph.space_after = Pt(12)
    
    return slide

def create_presentation(md_file, output_file):
    """Create PowerPoint presentation from markdown"""
    print(f"Reading markdown file: {md_file}")
    slides = parse_markdown_slides(md_file)
    print(f"Found {len(slides)} slides")
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Process each slide
    for i, slide_text in enumerate(slides, 1):
        title, content = parse_slide_content(slide_text)
        
        if not title:
            title = f"Slide {i}"
        
        print(f"Processing slide {i}: {title[:50]}...")
        add_slide(prs, title, content)
    
    # Save presentation
    prs.save(output_file)
    print(f"\n[SUCCESS] Presentation saved to: {output_file}")
    print(f"Total slides: {len(slides)}")

if __name__ == "__main__":
    md_file = "QuickAI_Presentation.md"
    output_file = "QuickAI_Presentation.pptx"
    
    try:
        create_presentation(md_file, output_file)
    except FileNotFoundError:
        print(f"[ERROR] {md_file} not found!")
        print("Make sure you're running this script from the QuickAI-Full-Stack directory")
    except ImportError:
        print("[ERROR] python-pptx not installed!")
        print("Install it using: pip install python-pptx")
    except Exception as e:
        print(f"[ERROR] {e}")

